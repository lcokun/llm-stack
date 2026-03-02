import csv
import io
import fitz # pymupdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from aiohttp import ClientSession
import chromadb

# --- Connect to ChromaDB ---

chroma = chromadb.HttpClient(host="chromadb", port=8000)
collection = chroma.get_or_create_collection("documents")

# --- File Readers
# Each one takes a file path and returns a single string of text

def read_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_xlsx(path):
    wb = load_workbook(path)
    lines = []
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            lines.append(" | ".join(str(c) for c in row if c is not None))
    return "\n".join(lines)

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text_frame.text)
    return "\n".join(text)

def read_csv_file(path):
    with open(path, "r") as f:
        reader = csv.reader(f)
        return "\n".join(" | ".join(row) for row in reader)

def read_text(path):
    with open(path, "r") as f:
        return f.read()

# Map extensions to headers
READERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".xlsx": read_xlsx,
    ".pptx": read_pptx,
    ".csv": read_csv_file,
    ".txt": read_text,
    ".md": read_text,
}

def read_file(path):
    """Pick the right header based on file extension."""
    ext = "." + path.rsplit(".", 1)[-1].lower()
    reader = READERS.get(ext)
    if not reader:
        raise ValueError(f"Unsupported file type: {ext}")
    return reader(path)

# --- Chunking ---
# Splits text into overlapping chunks so context isn't lost at boundaries

def chunk_text(text, chunk_size=500, overlap=50):
    """Split text into chunks of ~chunk_size words with overlap."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks

# --- Embedding  ----
# Sends text to Ollama's nomic-embed-text and gets back a vector

async def get_embedding(session, text, ollama_url):
    """Get embedding vector from Ollama."""
    payload = {"model": "nomic-embed-text", "input": text}
    async with session.post(f"{ollama_url}/api/embed", json=payload) as resp:
        data = await resp.json()
        return data["embeddings"][0]

# --- Ingest file ---
# The main function: read -> chunk -> embed -> store

async def ingest_file(path, ollama_url):
    """Read a file, chunk it, embed it, store in ChromaDB."""
    text = read_file(path)
    chunks = chunk_text(text)
    filename = path.rsplit("/", 1)[-1]

    async with ClientSession() as session:
        for i, chunk in enumerate(chunks):
            embedding = await get_embedding(session, chunk, ollama_url)
            doc_id = f"{filename}_chunk_{i}"

            collection.upsert(
                ids=[doc_id],
                embeddings=[embedding],
                documents=[chunk],
                metadatas=[{"source": filename, "chunk": i}],
            )

    return len(chunks)
